#!/usr/bin/env python3
"""
deck_to_pptx.py — convert an FGCZ-format slide deck (self-contained HTML) into an
editable .pptx.

Why this exists
---------------
The FGCZ presentation format (L1 `fgcz_presentation_format`) is authored as a
single self-contained HTML file: 1280x720 `.slide` divs, FGCZ palette, inline
CSS. That renders beautifully in a browser but cannot be handed to someone who
wants to edit it in PowerPoint. There is no headless browser on the FGCZ nodes,
so the usual HTML -> PDF -> image -> pptx route is unavailable.

Instead this script parses the deck's *semantic* markup and rebuilds each slide
from native PowerPoint shapes. The result is fully editable: real text boxes,
real rounded rectangles, real tables — not screenshots.

Geometry note: the HTML deck is 1280x720 px and a pptx 16:9 slide is
13.333 x 7.5 in, i.e. exactly 96 px per inch. So HTML pixel coordinates map 1:1
via PX(), and CSS px font sizes map to points as px * 0.75.

Supported block markup (keep new decks within this vocabulary so they convert):
    .title-slide  (.cover-title .cover-sub .cover-catch .cover-meta)
    .section-title .slide-title
    .two-col   > .col-card  (h3 + ul.bullet-list > li)
    .three-col > .col       (h3 + .sub + ul > li + .status)
    .four-col  > .col       (h3 + ul > li)
    .stackrows > .srow      (.srow-id .srow-q .srow-who)
    .flowrow   > .fbox      (b + span)   — .arrow separators are implied
    .stat-row  > .stat-card (h4 + .stat-line)
    ul.bullet-list          (standalone)
    .highlight-box[.orange|.pain|.proof]
    .key-point (.catch)
    table.spec
    .slide-footer

Usage:
    python3 deck_to_pptx.py deck_en.html [deck_ja.html ...]
    python3 deck_to_pptx.py --out /tmp deck_en.html

Requires: python-pptx, lxml.
"""

from __future__ import annotations

import argparse
import math
import os
import re
import sys

from lxml import html as LH
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palette ----
NAVY = RGBColor(0x00, 0x3C, 0x68)
NAVY_DARK = RGBColor(0x00, 0x14, 0x28)
CYAN = RGBColor(0x00, 0x99, 0xCC)
CYAN_LT = RGBColor(0x4D, 0xC9, 0xF6)
ORANGE = RGBColor(0xEA, 0x6B, 0x13)
GREY = RGBColor(0x54, 0x5E, 0x69)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE0, 0xE0, 0xE0)
SUBTLE = RGBColor(0xF9, 0xFA, 0xFB)
HDR_GREY = RGBColor(0xEE, 0xF1, 0xF6)

PAIN_BG, PAIN_BR = RGBColor(0xFE, 0xF2, 0xF2), RGBColor(0xEF, 0x44, 0x44)
PROOF_BG, PROOF_BR = RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0x10, 0xB9, 0x81)
SOL_BG = RGBColor(0xEF, 0xF6, 0xFF)
ORANGE_BG = RGBColor(0xFF, 0xF7, 0xED)
INSIGHT_BG, INSIGHT_BR = RGBColor(0xF5, 0xF3, 0xFF), RGBColor(0x8B, 0x5C, 0xF6)
NEW_BG = RGBColor(0xFF, 0xF4, 0xE2)

FONT_LATIN = "Arial"
FONT_CJK = "Noto Sans JP"
FONT_MONO = "Consolas"

CJK_RE = re.compile(r"[⺀-鿿＀-￯　-〿]")


def PX(px: float) -> Emu:
    """HTML pixel -> Emu (the deck is exactly 96 px per inch)."""
    return Inches(px / 96.0)


def SP(px: float) -> Pt:
    """CSS px font-size -> points."""
    return Pt(round(px * 0.75, 1))


def disp_len(s: str) -> float:
    """Display width in 'half-width units' (CJK counts double)."""
    return sum(2.0 if CJK_RE.match(c) else 1.0 for c in s)


# ------------------------------------------------------------ html helpers ----
def norm(s: str | None) -> str:
    return re.sub(r"\s+", " ", (s or "")).strip()


def cls(el) -> list[str]:
    return (el.get("class") or "").split()


def sel(el, want: str):
    """Direct-ish descendants carrying a class."""
    return [e for e in el.iter() if want in cls(e)]


def child_sel(el, want: str):
    return [e for e in el if want in cls(e)]


def runs_of(el) -> list[tuple[str, bool, bool]]:
    """Flatten an element into (text, bold, mono) runs, keeping <b>/<code>."""
    out: list[tuple[str, bool, bool]] = []

    def walk(node, bold, mono):
        if node.text:
            out.append((node.text, bold, mono))
        for ch in node:
            b = bold or ch.tag in ("b", "strong")
            m = mono or ch.tag == "code"
            walk(ch, b, m)
            if ch.tail:
                out.append((ch.tail, bold, mono))

    walk(el, el.tag in ("b", "strong"), el.tag == "code")
    merged: list[tuple[str, bool, bool]] = []
    for t, b, m in out:
        t = t.replace("\xa0", " ")
        if not t:
            continue
        if merged and merged[-1][1] == b and merged[-1][2] == m:
            merged[-1] = (merged[-1][0] + t, b, m)
        else:
            merged.append((t, b, m))
    if merged:
        merged[0] = (merged[0][0].lstrip(), merged[0][1], merged[0][2])
        merged[-1] = (merged[-1][0].rstrip(), merged[-1][1], merged[-1][2])
    return [(re.sub(r"\s+", " ", t), b, m) for t, b, m in merged if norm(t)]


def plain(el) -> str:
    return norm(el.text_content())


# ------------------------------------------------------------- pptx helpers ---
def _set_fonts(run, latin: str) -> None:
    """Set the latin face via python-pptx, and the East-Asian face via XML.

    python-pptx only exposes the latin typeface. Without an explicit <a:ea>,
    PowerPoint substitutes an arbitrary CJK font for Japanese runs, which
    changes the metrics the layout was computed for.
    """
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag, face in (("a:ea", FONT_CJK), ("a:cs", latin)):
        prefix, local = tag.split(":")
        qname = f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{local}"
        existing = rPr.find(qname)
        if existing is None:
            el = rPr.makeelement(qname, {"typeface": face})
            rPr.append(el)
        else:
            existing.set("typeface", face)


def set_text(tf, blocks, size_px, color=BLACK, bold_all=False, align=PP_ALIGN.LEFT,
             line_px=None, space_after_px=0, bullet=None):
    """blocks: list of run-lists (one per paragraph)."""
    tf.word_wrap = True
    first = True
    for runs in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if line_px:
            p.line_spacing = SP(line_px)
        if space_after_px:
            p.space_after = SP(space_after_px)
        prefix_done = False
        if not runs:
            runs = [("", False, False)]
        for text, bold, mono in runs:
            if bullet and not prefix_done:
                text = f"{bullet} {text.lstrip()}"
                prefix_done = True
            r = p.add_run()
            r.text = text
            f = r.font
            f.size = SP(size_px)
            f.bold = bold or bold_all
            f.color.rgb = NAVY if (bold and color is BLACK) else color
            _set_fonts(r, FONT_MONO if mono else FONT_LATIN)
    return tf


def textbox(slide, x, y, w, h, blocks, size_px, **kw):
    tb = slide.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    set_text(tf, blocks, size_px, **kw)
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False,
         radius=0.06):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        PX(x), PX(y), PX(w), PX(h),
    )
    if rounded:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = PX(12)
    tf.margin_top = tf.margin_bottom = PX(8)
    tf.word_wrap = True
    return shape


def est_h(runs_blocks, size_px, box_w_px, line_px=None, pad_px=18) -> float:
    """Estimate rendered height (px) of paragraph blocks in a box."""
    lh = line_px or size_px * 1.5
    # a half-width glyph is about 0.52 em wide in Arial/Noto at these sizes
    per_line = max(8.0, (box_w_px - 24) / (size_px * 0.52))
    lines = 0
    for runs in runs_blocks:
        t = "".join(r[0] for r in runs)
        lines += max(1, math.ceil(disp_len(t) / per_line))
    return lines * lh + pad_px


# ------------------------------------------------------------- slide render ---
class Deck:
    def __init__(self, path: str):
        self.path = path
        self.root = LH.parse(path).getroot()
        self.title = norm(self.root.findtext(".//title") or os.path.basename(path))
        self.prs = Presentation()
        self.prs.slide_width = PX(1280)
        self.prs.slide_height = PX(720)
        self.blank = self.prs.slide_layouts[6]
        self.warnings: list[str] = []

    # ---- chrome ----
    def _accent_bar(self, slide):
        # approximate the CSS 3-stop gradient with three bands
        for i, c in enumerate((NAVY, CYAN, CYAN_LT)):
            widths = (760, 300, 220)
            xs = (0, 760, 1060)
            rect(slide, xs[i], 0, widths[i], 8, fill=c)

    def _footer(self, slide, left, right):
        rect(slide, 0, 680, 1280, 1, fill=BORDER)
        textbox(slide, 60, 690, 400, 20, [[(left, True, False)]], 12,
                color=NAVY)
        tb = textbox(slide, 820, 690, 400, 20, [[(right, False, False)]], 12,
                     color=GREY, align=PP_ALIGN.RIGHT)
        return tb

    # ---- cover ----
    def _cover(self, slide, el):
        rect(slide, 0, 0, 1280, 720, fill=NAVY_DARK)
        rect(slide, 0, 0, 1280, 720, fill=NAVY)  # flat navy (gradient approx)
        t = sel(el, "cover-title")
        s = sel(el, "cover-sub")
        c = sel(el, "cover-catch")
        m = sel(el, "cover-meta")
        y = 210
        if t:
            textbox(slide, 90, y, 1100, 90, [[(plain(t[0]), True, False)]], 64,
                    color=WHITE, align=PP_ALIGN.CENTER)
            y += 110
        if s:
            textbox(slide, 190, y, 900, 60, [[(plain(s[0]), False, False)]], 20,
                    color=WHITE, align=PP_ALIGN.CENTER)
            y += 60
        if c:
            textbox(slide, 140, y + 20, 1000, 40, [[(plain(c[0]), True, False)]], 22,
                    color=CYAN_LT, align=PP_ALIGN.CENTER)
        if m:
            textbox(slide, 60, 630, 900, 24, [[(plain(m[0]), False, False)]], 13,
                    color=CYAN_LT)
        for r in range(3):
            for col in range(6):
                d = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, PX(1136 + col * 14), PX(632 + r * 14),
                    PX(8), PX(8))
                d.fill.solid()
                d.fill.fore_color.rgb = CYAN_LT
                d.line.fill.background()
                d.shadow.inherit = False

    # ---- blocks ----
    def _two_col(self, slide, el, x, y, w, h):
        cards = child_sel(el, "col-card")
        n = max(1, len(cards))
        gap = 26
        cw = (w - gap * (n - 1)) / n
        for i, card in enumerate(cards):
            cx = x + i * (cw + gap)
            box = rect(slide, cx, y, cw, h, fill=SUBTLE, line=BORDER,
                       rounded=True, radius=0.03)
            blocks = []
            head = card.find("h3")
            if head is not None:
                blocks.append(runs_of(head))
            for li in el_lis(card):
                blocks.append(runs_of(li))
            tf = box.text_frame
            tf.margin_left = tf.margin_right = PX(20)
            tf.margin_top = PX(14)
            set_text(tf, blocks, 14.5, line_px=21)
            # style the heading paragraph
            if head is not None and tf.paragraphs:
                for r in tf.paragraphs[0].runs:
                    r.font.size = SP(17)
                    r.font.bold = True
                    r.font.color.rgb = NAVY
                tf.paragraphs[0].space_after = SP(8)
            for p in tf.paragraphs[1:]:
                p.space_after = SP(3)
        return h

    def _n_col(self, slide, el, x, y, w, h, kind="three-col"):
        cols = [c for c in el if "col" in cls(c)]
        n = max(1, len(cols))
        gap = 18 if kind == "three-col" else 14
        cw = (w - gap * (n - 1)) / n
        tones = {
            "detect": (ORANGE_BG, ORANGE), "declare": (PROOF_BG, PROOF_BR),
            "defend": (SOL_BG, CYAN), "insight": (INSIGHT_BG, INSIGHT_BR),
            "results": (PROOF_BG, PROOF_BR), "learned": (INSIGHT_BG, INSIGHT_BR),
            "next": (SUBTLE, CYAN), "own": (ORANGE_BG, ORANGE),
        }
        for i, c in enumerate(cols):
            cx = x + i * (cw + gap)
            bg, top = SUBTLE, CYAN
            for k, (b, t) in tones.items():
                if k in cls(c):
                    bg, top = b, t
                    break
            rect(slide, cx, y, cw, h, fill=bg, line=None, rounded=True, radius=0.03)
            rect(slide, cx, y, cw, 4, fill=top)
            blocks, styles = [], []
            head = c.find("h3")
            if head is not None:
                blocks.append(runs_of(head)); styles.append("h")
            for s in child_sel(c, "sub"):
                blocks.append(runs_of(s)); styles.append("sub")
            for li in el_lis(c):
                blocks.append(runs_of(li)); styles.append("li")
            for s in child_sel(c, "status"):
                blocks.append(runs_of(s)); styles.append("status")
            tb = slide.shapes.add_textbox(PX(cx + 16), PX(y + 12),
                                          PX(cw - 32), PX(h - 20))
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            set_text(tf, blocks, 13, line_px=19)
            for p, st in zip(tf.paragraphs, styles):
                if st == "h":
                    for r in p.runs:
                        r.font.size = SP(17); r.font.bold = True
                        r.font.color.rgb = NAVY
                elif st == "sub":
                    for r in p.runs:
                        r.font.size = SP(11); r.font.color.rgb = GREY
                    p.space_after = SP(8)
                elif st == "status":
                    for r in p.runs:
                        r.font.size = SP(11); r.font.italic = True
                        r.font.color.rgb = GREY
                    p.space_before = SP(8)
                else:
                    p.space_after = SP(3)
                    if p.runs:
                        p.runs[0].text = "• " + p.runs[0].text
        return h

    def _stackrows(self, slide, el, x, y, w, h):
        rows = child_sel(el, "srow")
        n = max(1, len(rows))
        rh = h / n
        rect(slide, x, y, w, h, fill=None, line=BORDER)
        idw, qw = 64, (w - 64) * 0.53
        for i, r in enumerate(rows):
            ry = y + i * rh
            is_new = "new" in cls(r)
            rect(slide, x, ry, idw, rh, fill=NEW_BG if is_new else HDR_GREY,
                 line=BORDER)
            if is_new:
                rect(slide, x + idw, ry, qw, rh,
                     fill=RGBColor(0xFF, 0xFA, 0xF2), line=None)
            for want, cx, cw_, sz, col in (
                ("srow-id", x, idw, 13, ORANGE if is_new else NAVY),
                ("srow-q", x + idw, qw, 13.5, BLACK),
                ("srow-who", x + idw + qw, w - idw - qw, 12.5, GREY),
            ):
                cells = child_sel(r, want)
                if not cells:
                    continue
                al = PP_ALIGN.CENTER if want == "srow-id" else PP_ALIGN.LEFT
                tb = slide.shapes.add_textbox(PX(cx + 10), PX(ry + 5),
                                              PX(cw_ - 20), PX(rh - 8))
                tf = tb.text_frame
                tf.margin_left = tf.margin_right = 0
                tf.margin_top = tf.margin_bottom = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                set_text(tf, [runs_of(cells[0])], sz, color=col, align=al,
                         line_px=sz * 1.45,
                         bold_all=(want == "srow-id"))
            if i:
                rect(slide, x, ry, w, 1, fill=BORDER)
        return h

    def _flowrow(self, slide, el, x, y, w, h):
        boxes = child_sel(el, "fbox")
        n = max(1, len(boxes))
        arrow_w, gap = 26, 8
        bw = (w - (n - 1) * (arrow_w + gap * 2)) / n
        for i, b in enumerate(boxes):
            bx = x + i * (bw + arrow_w + gap * 2)
            fill, line = SUBTLE, NAVY
            if "new" in cls(b):
                fill, line = ORANGE_BG, ORANGE
            elif "ok" in cls(b):
                fill, line = PROOF_BG, PROOF_BR
            elif "no" in cls(b):
                fill, line = PAIN_BG, PAIN_BR
            rect(slide, bx, y, bw, h, fill=fill, line=line, line_w=1.5,
                 rounded=True, radius=0.08)
            blocks, styles = [], []
            hb = b.find("b")
            if hb is not None:
                blocks.append(runs_of(hb)); styles.append("h")
            for s in b.findall("span"):
                blocks.append(runs_of(s)); styles.append("d")
            tb = slide.shapes.add_textbox(PX(bx + 13), PX(y + 10),
                                          PX(bw - 26), PX(h - 18))
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            set_text(tf, blocks, 12, line_px=17)
            for p, st in zip(tf.paragraphs, styles):
                if st == "h":
                    for r in p.runs:
                        r.font.size = SP(14); r.font.bold = True
                        r.font.color.rgb = NAVY
                    p.space_after = SP(3)
                else:
                    for r in p.runs:
                        r.font.color.rgb = GREY
            if i < n - 1:
                ax = bx + bw + gap
                ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, PX(ax),
                                            PX(y + h / 2 - 9), PX(arrow_w), PX(18))
                ar.fill.solid(); ar.fill.fore_color.rgb = ORANGE
                ar.line.fill.background(); ar.shadow.inherit = False
        return h

    def _stat_row(self, slide, el, x, y, w, h):
        cards = child_sel(el, "stat-card")
        n = max(1, len(cards))
        gap = 18
        cw = (w - gap * (n - 1)) / n
        for i, c in enumerate(cards):
            cx = x + i * (cw + gap)
            accent = PAIN_BR if "before" in cls(c) else PROOF_BR
            rect(slide, cx, y, cw, h, fill=WHITE, line=BORDER, rounded=True,
                 radius=0.03)
            rect(slide, cx, y, 4, h, fill=accent)
            blocks, styles = [], []
            hd = c.find("h4")
            if hd is not None:
                blocks.append(runs_of(hd)); styles.append("h")
            for ln in child_sel(c, "stat-line"):
                blocks.append(runs_of(ln)); styles.append("l")
            tb = slide.shapes.add_textbox(PX(cx + 18), PX(y + 12),
                                          PX(cw - 36), PX(h - 20))
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            set_text(tf, blocks, 14, line_px=20)
            for p, st in zip(tf.paragraphs, styles):
                if st == "h":
                    for r in p.runs:
                        r.font.size = SP(11); r.font.color.rgb = GREY
                        r.font.bold = True
                    p.space_after = SP(6)
                else:
                    for r in p.runs:
                        if r.font.bold:
                            r.font.size = SP(17)
                            r.font.color.rgb = accent
        return h

    def _highlight(self, slide, el, x, y, w, h):
        c = cls(el)
        bg, br = SOL_BG, CYAN
        if "orange" in c:
            bg, br = ORANGE_BG, ORANGE
        elif "pain" in c:
            bg, br = PAIN_BG, PAIN_BR
        elif "proof" in c:
            bg, br = PROOF_BG, PROOF_BR
        rect(slide, x, y, w, h, fill=bg, line=None)
        rect(slide, x, y, 4, h, fill=br)
        tb = slide.shapes.add_textbox(PX(x + 16), PX(y + 8), PX(w - 32), PX(h - 14))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        set_text(tf, [runs_of(el)], 14.5, line_px=21)
        return h

    def _key_point(self, slide, el, x, y, w, h):
        rect(slide, x, y, w, h, fill=NAVY_DARK, line=None, rounded=True,
             radius=0.04)
        catch = child_sel(el, "catch")
        tb = slide.shapes.add_textbox(PX(x + 24), PX(y + 10), PX(w - 48), PX(h - 20))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if catch:
            set_text(tf, [[(plain(catch[0]), True, False)]], 28, color=ORANGE,
                     align=PP_ALIGN.CENTER)
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.italic = True
        else:
            paras = [runs_of(el)]
            set_text(tf, paras, 20, color=WHITE, align=PP_ALIGN.CENTER,
                     line_px=30)
            for p in tf.paragraphs:
                for r in p.runs:
                    if r.font.bold:
                        r.font.color.rgb = CYAN_LT
        return h

    def _table(self, slide, el, x, y, w, h):
        trs = [tr for tr in el.iter("tr")]
        if not trs:
            return 0
        ncols = max(len(list(tr)) for tr in trs)
        gf = slide.shapes.add_table(len(trs), ncols, PX(x), PX(y), PX(w), PX(h))
        tbl = gf.table
        widths = []
        for cell in trs[0]:
            m = re.search(r"width:\s*(\d+)px", cell.get("style") or "")
            widths.append(int(m.group(1)) if m else None)
        known = sum(v for v in widths if v)
        free = [i for i, v in enumerate(widths) if not v]
        for i, v in enumerate(widths):
            if v:
                tbl.columns[i].width = PX(v)
            elif free:
                tbl.columns[i].width = PX((w - known) / len(free))
        for ri, tr in enumerate(trs):
            cells = list(tr)
            is_head = cells and cells[0].tag == "th"
            tbl.rows[ri].height = PX(max(20, (h - 26) / max(1, len(trs) - 1)))
            for ci in range(ncols):
                cell = tbl.cell(ri, ci)
                cell.margin_left = cell.margin_right = PX(8)
                cell.margin_top = cell.margin_bottom = PX(3)
                cell.fill.solid()
                cell.fill.fore_color.rgb = HDR_GREY if is_head else WHITE
                txt = runs_of(cells[ci]) if ci < len(cells) else []
                tf = cell.text_frame
                tf.word_wrap = True
                set_text(tf, [txt], 12 if not is_head else 11.5,
                         color=NAVY if is_head else BLACK,
                         bold_all=is_head, line_px=17)
        return h

    # ---- one content slide ----
    def _content(self, slide, el):
        self._accent_bar(slide)
        content = sel(el, "slide-content")
        content = content[0] if content else el
        x, w = 60, 1160
        y = 44
        for st in child_sel(content, "section-title"):
            textbox(slide, x, y, w, 22, [[(plain(st).upper(), True, False)]], 15,
                    color=CYAN)
            y += 30
        for tt in child_sel(content, "slide-title"):
            txt = plain(tt)
            lines = max(1, math.ceil(disp_len(txt) / 46))
            textbox(slide, x, y, w, 44 * lines,
                    [[(txt, True, False)]], 34, color=NAVY, line_px=42)
            y += 44 * lines + 8

        # collect blocks with their natural height
        kinds = (("two-col", "two"), ("three-col", "three"), ("four-col", "four"),
                 ("stackrows", "stack"), ("flowrow", "flow"), ("stat-row", "stat"),
                 ("highlight-box", "hl"), ("key-point", "key"))
        blocks = []
        for ch in content:
            c = cls(ch)
            if "section-title" in c or "slide-title" in c:
                continue
            kind = next((k for cname, k in kinds if cname in c), None)
            if kind is None:
                kind = {"table": "table", "ul": "ul"}.get(ch.tag)
            if kind is None:
                continue
            blocks.append([kind, ch, nat_h(kind, ch, w)])

        # fit: distribute slack to the blocks that benefit, shrink if over
        avail = 644 - y
        gap = 12
        total_gap = gap * max(0, len(blocks) - 1)
        space = avail - total_gap
        nsum = sum(b[2] for b in blocks) or 1
        if nsum < space:
            grow = [b for b in blocks if b[0] in
                    ("two", "three", "four", "stack", "table", "stat")] or blocks
            gsum = sum(b[2] for b in grow)
            slack = space - nsum
            for b in grow:
                b[2] += slack * b[2] / gsum
        elif nsum > space:
            for b in blocks:
                b[2] *= space / nsum
            self.warnings.append(
                f"content scaled to {space / nsum:.0%} on a slide "
                f"(title: {plain(content.find('.//div')) [:40]!r})")

        for kind, ch, hh in blocks:
            if kind == "two":
                self._two_col(slide, ch, x, y, w, hh)
            elif kind in ("three", "four"):
                self._n_col(slide, ch, x, y, w, hh,
                            "three-col" if kind == "three" else "four-col")
            elif kind == "stack":
                self._stackrows(slide, ch, x, y, w, hh)
            elif kind == "flow":
                self._flowrow(slide, ch, x, y, w, hh)
            elif kind == "stat":
                self._stat_row(slide, ch, x, y, w, hh)
            elif kind == "hl":
                self._highlight(slide, ch, x, y, w, hh)
            elif kind == "key":
                self._key_point(slide, ch, x, y, w, hh)
            elif kind == "table":
                self._table(slide, ch, x, y, w, hh)
            elif kind == "ul":
                textbox(slide, x, y, w, hh,
                        [runs_of(li) for li in ch.iter("li")], 14.5,
                        line_px=21, bullet="–")
            y += hh + gap

        foot = sel(el, "slide-footer")
        if foot:
            spans = list(foot[0])
            left = plain(spans[0]) if spans else "FGCZ · GENOMICS"
            right = plain(spans[1]) if len(spans) > 1 else ""
            self._footer(slide, left, right)

    # ---- driver ----
    def build(self) -> Presentation:
        for el in self.root.iter():
            if "slide" not in cls(el):
                continue
            slide = self.prs.slides.add_slide(self.blank)
            if "title-slide" in cls(el):
                self._cover(slide, el)
            else:
                self._content(slide, el)
        return self.prs


def el_lis(el):
    return list(el.iter("li"))


def nat_h(kind: str, el, w: float) -> float:
    """Natural (text-driven) height in px for one content block."""
    if kind in ("two", "stat"):
        cards = [c for c in el]
        n = max(1, len(cards))
        cw = (w - 26 * (n - 1)) / n
        best = 0.0
        for c in cards:
            blocks = []
            head = c.find("h3") if kind == "two" else c.find("h4")
            if head is not None:
                blocks.append(runs_of(head))
            rows = el_lis(c) if kind == "two" else child_sel(c, "stat-line")
            blocks += [runs_of(r) for r in rows]
            sz, lp = (14.5, 21) if kind == "two" else (14, 20)
            best = max(best, est_h(blocks, sz, cw, lp, 40))
        return max(90.0, best)
    if kind in ("three", "four"):
        cols = [c for c in el if "col" in cls(c)]
        n = max(1, len(cols))
        gap = 18 if kind == "three" else 14
        cw = (w - gap * (n - 1)) / n
        best = 0.0
        for c in cols:
            blocks = []
            head = c.find("h3")
            if head is not None:
                blocks.append(runs_of(head))
            blocks += [runs_of(s) for s in child_sel(c, "sub")]
            blocks += [runs_of(li) for li in el_lis(c)]
            blocks += [runs_of(s) for s in child_sel(c, "status")]
            best = max(best, est_h(blocks, 13, cw, 19, 46))
        return max(120.0, best)
    if kind == "stack":
        rows = child_sel(el, "srow")
        total = 0.0
        idw = 64
        qw = (w - idw) * 0.53
        for r in rows:
            hs = [30.0]
            for want, cw in (("srow-q", qw), ("srow-who", w - idw - qw)):
                for cell in child_sel(r, want):
                    sz = 13.5 if want == "srow-q" else 12.5
                    hs.append(est_h([runs_of(cell)], sz, cw, sz * 1.45, 14))
            total += max(hs)
        return total
    if kind == "flow":
        boxes = child_sel(el, "fbox")
        n = max(1, len(boxes))
        bw = (w - (n - 1) * 42) / n
        best = 0.0
        for b in boxes:
            blocks = []
            hb = b.find("b")
            if hb is not None:
                blocks.append(runs_of(hb))
            blocks += [runs_of(s) for s in b.findall("span")]
            best = max(best, est_h(blocks, 12, bw, 17, 30))
        return max(90.0, best)
    if kind == "hl":
        return max(48.0, est_h([runs_of(el)], 14.5, w, 21, 20))
    if kind == "key":
        if child_sel(el, "catch"):
            return 74.0
        return max(70.0, est_h([runs_of(el)], 20, w, 30, 26))
    if kind == "table":
        return 26.0 + max(1, len(list(el.iter("tr"))) - 1) * 25.0
    if kind == "ul":
        return est_h([runs_of(li) for li in el_lis(el)], 14.5, w, 21, 10)
    return 60.0


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("decks", nargs="+", help="deck HTML file(s)")
    ap.add_argument("--out", default=None, help="output directory")
    args = ap.parse_args()

    rc = 0
    for path in args.decks:
        if not os.path.isfile(path):
            print(f"!! not found: {path}", file=sys.stderr)
            rc = 1
            continue
        deck = Deck(path)
        prs = deck.build()
        outdir = args.out or os.path.dirname(os.path.abspath(path))
        out = os.path.join(outdir,
                           os.path.splitext(os.path.basename(path))[0] + ".pptx")
        prs.save(out)
        print(f"{os.path.basename(path)} -> {os.path.basename(out)}  "
              f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
              f"{os.path.getsize(out) // 1024} KB)")
        for wmsg in deck.warnings:
            print(f"   warn: {wmsg}")
    return rc


if __name__ == "__main__":
    sys.exit(main())
